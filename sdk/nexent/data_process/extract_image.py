import os
import base64
import hashlib
import tempfile
import subprocess
from typing import List, Dict, Any, Optional
import zipfile
from xml.etree import ElementTree

from pptx import Presentation

from .base import FileProcessor

from unstructured_inference.logger import logger
from unstructured_inference.models import tables
from unstructured.partition.auto import partition


tables_agent = tables.tables_agent
TABLE_TRANSFORMER_MODEL_PATH = ""

def custom_load_table_model():
    """Loads the Table agent."""

    if getattr(tables_agent, "model", None) is None:
        with tables_agent._lock:
            if getattr(tables_agent, "model", None) is None:
                logger.info("Loading the Table agent ...")
                print("path234: ", TABLE_TRANSFORMER_MODEL_PATH)
                tables_agent.initialize(TABLE_TRANSFORMER_MODEL_PATH)

    return

tables.load_agent = lambda: custom_load_table_model()


class UniversalImageExtractor(FileProcessor):
    """
    Multi-format image extractor for PDF, PPT, Excel, and Word.
    Uses LibreOffice for conversion when needed and reuses PDF extraction logic.
    """

    @staticmethod
    def _hash(data: bytes) -> str:
        # Use a modern hash for safe, collision-resistant de-duplication.
        return hashlib.sha256(data).hexdigest()

    @staticmethod
    def _openxml_namespace_maps() -> Dict[str, str]:
        return {
            "xdr": "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing",  # NOSONAR
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",  # NOSONAR
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",  # NOSONAR
        }


    def _write_temp_file(self, data: bytes, suffix: str) -> str:
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(data)
        tmp.close()
        return tmp.name

    @staticmethod
    def detect_image_format(image_bytes: bytes) -> str:
        if image_bytes.startswith(b"\x89PNG"):
            return "png"
        elif image_bytes.startswith(b"\xFF\xD8\xFF"):
            return "jpg"
        else:
            return "png"


    def _convert_file(self, input_path: str, target_format: str) -> str:
    
        """
        Convert a file to the target format using LibreOffice.

        Args:
            input_path: Source file path.
            target_format: Target format, e.g. "pdf", "pptx", "xlsx".

        Returns:
            Output file path.
        """
        out_dir = os.path.dirname(input_path)

        cmd = [
            "soffice",
            "--headless",
            "--invisible",  # Ensure fully headless conversion.
            "--convert-to", f"{target_format}",
            input_path,
            "--outdir", out_dir
        ]

        try:
            subprocess.run(
                cmd,
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=60  # Prevent hanging conversions.
            )

            base_name = os.path.splitext(input_path)[0]
            new_suffix = f".{target_format}"
            output_path = base_name + new_suffix

            if os.path.exists(output_path):
                return output_path
            else:
                raise FileNotFoundError(
                    f"Conversion failed: Output file {output_path} not found.")

        except subprocess.CalledProcessError as e:
            raise RuntimeError(
                f"LibreOffice conversion failed for {input_path}: {e}")
        except subprocess.TimeoutExpired:
            raise RuntimeError(
                f"LibreOffice conversion timed out for {input_path}")


    def _extract_pdf(self, pdf_path: str, **params) -> List[Dict]:
        table_transformer_model_path = params.get("table_transformer_model_path")
        unstructured_default_model_initialize_params_json_path = params.get(
            "unstructured_default_model_initialize_params_json_path"
        )
        if not table_transformer_model_path or not unstructured_default_model_initialize_params_json_path:
            return []
        global TABLE_TRANSFORMER_MODEL_PATH
        TABLE_TRANSFORMER_MODEL_PATH = table_transformer_model_path

        results = []
        seen = set()

        elements = partition(
            filename=pdf_path,
            strategy="hi_res",
            extract_images_in_pdf=True,
            extract_image_block_to_payload=True,
        )

        for el in elements:
            b64 = getattr(el.metadata, "image_base64", None)
            if not b64:
                continue

            img_bytes = base64.b64decode(b64)
            h = self._hash(img_bytes)
            if h in seen:
                continue
            seen.add(h)

            coords = getattr(el.metadata, "coordinates", None)
            coord_dict = None

            if coords and hasattr(coords, 'points') and coords.points:
                pts = coords.points  # tuple of (x,y)
                xs = [p[0] for p in pts]
                ys = [p[1] for p in pts]
                coord_dict = {
                    "x1": min(xs),
                    "y1": min(ys),
                    "x2": max(xs),
                    "y2": max(ys),
                }

            page_num = getattr(el.metadata, "page_number", None)

            results.append({
                "position": {
                    "page_number": page_num,
                    "coordinates": coord_dict
                },
                "image_format": self.detect_image_format(img_bytes),
                "image_bytes": img_bytes
            })

        return results


    def _excel_sheet_files(self, z: zipfile.ZipFile) -> List[str]:
        return [f for f in z.namelist() if f.startswith("xl/worksheets/sheet")]


    def _excel_drawing_file(self, z: zipfile.ZipFile, sheet_file: str) -> Optional[str]:
        sheet_xml = ElementTree.fromstring(z.read(sheet_file))
        drawing = sheet_xml.find(
            ".//{https://schemas.openxmlformats.org/spreadsheetml/2006/main}drawing")
        if drawing is None:
            drawing = sheet_xml.find(
                ".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}drawing")
        if drawing is None:
            return None

        rel_id = drawing.get(
            "{https://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rel_id is None:
            rel_id = drawing.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        rel_path = sheet_file.replace("worksheets", "worksheets/_rels") + ".rels"
        if rel_path not in z.namelist():
            return None

        rel_xml = ElementTree.fromstring(z.read(rel_path))
        for rel in rel_xml:
            if rel.get("Id") == rel_id:
                return "xl/" + rel.get("Target").replace("../", "")

        return None


    def _excel_rel_map(self, z: zipfile.ZipFile, drawing_file: str) -> Optional[Dict[str, str]]:
        rel_file = drawing_file.replace("drawings/", "drawings/_rels/") + ".rels"
        if rel_file not in z.namelist():
            return None

        rel_root = ElementTree.fromstring(z.read(rel_file))
        return {
            rel.get("Id"): "xl/" + rel.get("Target").replace("../", "")
            for rel in rel_root
        }


    def _excel_anchors(self, z: zipfile.ZipFile, drawing_file: str, ns: Dict[str, str]) -> List[Any]:
        drawing_root = ElementTree.fromstring(z.read(drawing_file))
        return drawing_root.findall(".//xdr:twoCellAnchor", ns) + \
            drawing_root.findall(".//xdr:oneCellAnchor", ns)


    def _excel_anchor_coords(self, anchor: Any, ns: Dict[str, str]) -> Optional[Dict[str, int]]:
        from_node = anchor.find("xdr:from", ns)
        if from_node is None:
            return None

        row1 = int(from_node.find("xdr:row", ns).text) + 1
        col1 = int(from_node.find("xdr:col", ns).text) + 1

        to_node = anchor.find("xdr:to", ns)
        if to_node is not None:
            row2 = int(to_node.find("xdr:row", ns).text) + 1
            col2 = int(to_node.find("xdr:col", ns).text) + 1
        else:
            row2, col2 = row1, col1

        return {"row1": row1, "col1": col1, "row2": row2, "col2": col2}


    def _excel_anchor_embed_id(self, anchor: Any, ns: Dict[str, str]) -> Optional[str]:
        blip = anchor.find(".//a:blip", ns)
        if blip is None:
            return None

        embed_id = blip.get(
            "{https://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
        if embed_id is None:
            embed_id = blip.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
        return embed_id


    def _extract_excel_anchors(
        self,
        z: zipfile.ZipFile,
        anchors: List[Any],
        rel_map: Dict[str, str],
        sheet_name: str,
        ns: Dict[str, str],
        seen: set,
    ) -> List[Dict[str, Any]]:
        results = []
        for anchor in anchors:
            coords = self._excel_anchor_coords(anchor, ns)
            if coords is None:
                continue

            embed_rel_id = self._excel_anchor_embed_id(anchor, ns)
            if not embed_rel_id:
                continue

            target = rel_map.get(embed_rel_id)
            if not target:
                continue

            img_bytes = z.read(target)
            h = self._hash(img_bytes)
            if h in seen:
                continue
            seen.add(h)

            results.append({
                "position": {
                    "sheet_name": sheet_name,
                    "coordinates": {
                        "x1": coords["col1"],
                        "x2": coords["col2"],
                        "y1": coords["row1"],
                        "y2": coords["row2"]
                    }
                },
                "image_format": self.detect_image_format(img_bytes),
                "image_bytes": img_bytes
            })

        return results


    def _extract_excel_sheet(
        self,
        z: zipfile.ZipFile,
        sheet_file: str,
        ns: Dict[str, str],
        seen: set,
    ) -> List[Dict[str, Any]]:
        drawing_file = self._excel_drawing_file(z, sheet_file)
        if drawing_file is None:
            return []

        rel_map = self._excel_rel_map(z, drawing_file)
        if not rel_map:
            return []

        anchors = self._excel_anchors(z, drawing_file, ns)
        sheet_name = os.path.basename(sheet_file)

        return self._extract_excel_anchors(z, anchors, rel_map, sheet_name, ns, seen)


    def _extract_excel(self, xlsx_path):
        results = []
        seen = set()

        with zipfile.ZipFile(xlsx_path) as z:
            sheet_files = self._excel_sheet_files(z)

            ns = self._openxml_namespace_maps()
            for sheet_file in sheet_files:
                results.extend(self._extract_excel_sheet(z, sheet_file, ns, seen))

        return results


    def _extract_pptx(self, pptx_path: str, **params) -> List[Dict]:
        if Presentation is None:
            raise RuntimeError("python-pptx is required to extract images from PPTX files.")
        prs = Presentation(pptx_path)
        results = []
        seen = set()
        emu_per_inch = params.get("emu_per_inch", 914400)
        dpi = params.get("dpi", 96)
        
        def _emu_to_px(emu: int, emu_per_inch: int, dpi: int) -> int:
            return int((emu / emu_per_inch) * dpi)
        

        slide_w = _emu_to_px(prs.slide_width, emu_per_inch, dpi)
        slide_h = _emu_to_px(prs.slide_height, emu_per_inch, dpi)

        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, "image"):
                    continue

                img_bytes = shape.image.blob
                h = self._hash(img_bytes)
                if h in seen:
                    continue
                seen.add(h)

                x = _emu_to_px(shape.left, emu_per_inch, dpi)
                y = _emu_to_px(shape.top, emu_per_inch, dpi)
                w = _emu_to_px(shape.width, emu_per_inch, dpi)
                h_px = _emu_to_px(shape.height, emu_per_inch, dpi)

                results.append({
                    "position": {
                        "page_number": slide_index + 1,
                        "coordinates": {
                            "x1": x,
                            "y1": y,
                            "x2": x + w,
                            "y2": y + h_px,
                            "slide_width": slide_w,
                            "slide_height": slide_h,
                        },
                    },
                    "image_format": self.detect_image_format(img_bytes),
                    "image_bytes": img_bytes
                })

        return results


    def process_file(self, file_bytes: bytes, chunking_strategy: str, filename: str, **params) -> List[Dict[str, Any]]:
        suffix = os.path.splitext(filename)[1].lower()
        temp_path = self._write_temp_file(file_bytes, suffix)
        converted_path = None

        try:
            direct_extractors = {
                ".xlsx": lambda: self._extract_excel(temp_path),
                ".pptx": lambda: self._extract_pptx(temp_path, **params),
                ".pdf": lambda: self._extract_pdf(temp_path, **params),
            }
            if suffix in direct_extractors:
                return direct_extractors[suffix]()

            conversions = {
                ".xls": ("xlsx", lambda path: self._extract_excel(path)),
                ".ppt": ("pptx", lambda path: self._extract_pptx(path, **params)),
                ".docx": ("pdf", lambda path: self._extract_pdf(path, **params)),
                ".doc": ("pdf", lambda path: self._extract_pdf(path, **params)),
            }
            if suffix in conversions:
                target_format, extractor = conversions[suffix]
                converted_path = self._convert_file(temp_path, target_format)
                return extractor(converted_path)

            return []

        finally:
            files_to_clean = [temp_path]
            if converted_path and os.path.exists(converted_path):
                files_to_clean.append(converted_path)

            base = os.path.splitext(temp_path)[0]
            for ext in [".docx", ".pptx", ".xlsx", ".pdf"]:
                potential_file = base + ext
                if potential_file != converted_path and potential_file != temp_path:
                    files_to_clean.append(potential_file)

            for f_path in files_to_clean:
                if f_path and os.path.exists(f_path):
                    try:
                        os.remove(f_path)
                    except Exception:
                        pass
